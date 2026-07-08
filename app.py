import streamlit as st
import pdfplumber
import pandas as pd
import re
import logging
import os
import struct
import zlib
from docx import Document
from pptx import Presentation
import olefile

# 지저분한 로그 메시지 숨기기
logging.getLogger("pdfminer").setLevel(logging.ERROR)

st.set_page_config(page_title="만능 문서 마크다운 변환기", page_icon="📝", layout="wide")

st.title("📑 만능 문서 마크다운 자동 변환기")
st.markdown("""
- **지원 형식:** PDF, Word(.docx), PowerPoint(.pptx), 한글(.hwp)
""")

allowed_types = ["pdf", "docx", "pptx", "hwp"]
uploaded_files = st.file_uploader("변환할 파일들을 선택하세요", type=allowed_types, accept_multiple_files=True)

def clean_text(text):
    if not text: return ""
    text = re.sub(r'[ ]+', ' ', text)
    return text.strip()

def extract_pdf(file):
    full_md = []
    with pdfplumber.open(file) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text(x_tolerance=2, y_tolerance=2)
            if page_text:
                full_md.append(f"## Page {i+1}\n\n{clean_text(page_text)}\n")
    return "\n".join(full_md)

def extract_docx(file):
    doc = Document(file)
    full_md = ["## Word Document 본문\n"]
    for para in doc.paragraphs:
        if para.text.strip():
            full_md.append(clean_text(para.text))
    return "\n\n".join(full_md)

def extract_pptx(file):
    prs = Presentation(file)
    full_md = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_md = []
                for row_idx, row in enumerate(table.rows):
                    row_text = [clean_text(cell.text) for cell in row.cells]
                    table_md.append("| " + " | ".join(row_text) + " |")
                    if row_idx == 0:
                        separator = "|" + "|".join(["---"] * len(row.cells)) + "|"
                        table_md.append(separator)
                if table_md:
                    slide_text.append("\n" + "\n".join(table_md) + "\n")
            elif hasattr(shape, "text") and shape.text.strip():
                slide_text.append(clean_text(shape.text))
        if slide_text:
            full_md.append(f"## Slide {i+1}\n\n" + "\n\n".join(slide_text) + "\n")
    return "\n".join(full_md)

def extract_hwp(file):
    try:
        ole_obj = olefile.OleFileIO(file)
        dirs = ole_obj.listdir()
        sections = [d for d in dirs if d[0] == 'BodyText']
        full_md = ["## 한글 문서(HWP) 본문\n"]
        for section in sections:
            stream = ole_obj.openstream(section)
            data = stream.read()
            try: unpacked_data = zlib.decompress(data, -15)
            except: unpacked_data = data
            i = 0
            section_text = []
            while i < len(unpacked_data):
                if i + 4 > len(unpacked_data): break
                header = struct.unpack('<I', unpacked_data[i:i+4])[0]
                tag_id = header & 0x3FF
                size = (header >> 20) & 0xFFF
                i += 4
                if size == 0xFFF:
                    if i + 4 > len(unpacked_data): break
                    size = struct.unpack('<I', unpacked_data[i:i+4])[0]
                    i += 4
                if tag_id == 67:
                    record_data = unpacked_data[i:i+size]
                    text = record_data.decode('utf-16-le', errors='ignore')
                    cleaned = "".join([c for c in text if ord(c) >= 32 or c in ('\n', '\t')])
                    if cleaned.strip(): section_text.append(cleaned)
                i += size
            if section_text:
                full_md.append("\n\n".join(section_text))
        return "\n\n".join(full_md)
    except Exception as e:
        return f"HWP 추출 중 오류 발생: {str(e)}"

if uploaded_files:
    st.subheader("✅ 변환 리스트")
    for uploaded_file in uploaded_files:
        col1, col2, col3 = st.columns([5, 2, 2])
        ext = uploaded_file.name.split('.')[-1].lower()
        with col1:
            st.write(f"📄 **{uploaded_file.name}** ({ext.upper()})")
        with col2:
            status_placeholder = st.empty()
            status_placeholder.text("⏳ 처리 중...")
        try:
            md_result = ""
            if ext == "pdf": md_result = extract_pdf(uploaded_file)
            elif ext == "docx": md_result = extract_docx(uploaded_file)
            elif ext == "pptx": md_result = extract_pptx(uploaded_file)
            elif ext == "hwp": md_result = extract_hwp(uploaded_file)

            if md_result.strip():
                status_placeholder.text("✨ 변환 완료!")
                with col3:
                    base_name = os.path.splitext(uploaded_file.name)[0]
                    st.download_button(
                        label="📥 MD 다운로드",
                        data=md_result,
                        file_name=f"{base_name}.md",
                        mime="text/markdown",
                        key=f"btn_{uploaded_file.name}"
                    )
            else:
                status_placeholder.text("⚠️ 추출된 텍스트 없음")
        except Exception as e:
            status_placeholder.text("❌ 오류 발생")
            st.error(f"에러 내용: {e}")
    st.divider()
